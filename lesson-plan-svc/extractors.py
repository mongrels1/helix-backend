"""
extractors.py — Turn any accepted resource into normalized text, and
auto-classify it (pacing / ilp / perf_task / standards / slides / roster / other).

Supported: .docx .pdf .xlsx/.xls .csv .pptx .txt/.md .png/.jpg (OCR if tesseract
present) and web URLs. Each extractor is defensive: a failure on one file never
breaks the batch — it returns an error note instead.
"""
from __future__ import annotations
import os, re, io

def extract_text(path: str, filename: str = "") -> str:
    name = (filename or os.path.basename(path)).lower()
    ext = os.path.splitext(name)[1]
    try:
        if ext == ".docx":
            return _docx(path)
        if ext == ".pdf":
            return _pdf(path)
        if ext in (".xlsx", ".xlsm", ".xls"):
            return _xlsx(path)
        if ext == ".csv":
            return _csv(path)
        if ext == ".pptx":
            return _pptx(path)
        if ext in (".txt", ".md"):
            return _plain(path)
        if ext in (".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"):
            return _ocr(path)
        return _plain(path)
    except Exception as e:
        return f"[could not extract {name}: {e}]"

def extract_url(url: str) -> str:
    try:
        import requests
        from bs4 import BeautifulSoup
        r = requests.get(url, timeout=15, headers={"User-Agent": "EdKairos-LPM/0.1"})
        soup = BeautifulSoup(r.text, "html.parser")
        for t in soup(["script", "style", "nav", "footer", "header"]):
            t.decompose()
        text = re.sub(r"\n{3,}", "\n\n", soup.get_text("\n"))
        return f"[web:{url}]\n{text.strip()[:8000]}"
    except Exception as e:
        return f"[could not fetch {url}: {e}]"

# ---- per-format ----
def _docx(path):
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for t in d.tables:
        for r in t.rows:
            cells = [c.text.strip() for c in r.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)

def _pdf(path):
    import pdfplumber
    out = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages[:40]:
            out.append(page.extract_text() or "")
            for tbl in (page.extract_tables() or []):
                for row in tbl:
                    out.append(" | ".join(c or "" for c in row))
    return "\n".join(out)

def _xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"# sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            vals = [str(c) for c in row if c is not None]
            if vals:
                out.append(" | ".join(vals))
    return "\n".join(out[:2000])

def _csv(path):
    import csv
    out = []
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        for row in csv.reader(f):
            if any(row):
                out.append(" | ".join(row))
    return "\n".join(out[:2000])

def _pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"# slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs)
                    if t.strip():
                        out.append(t)
    return "\n".join(out)

def _plain(path):
    with open(path, encoding="utf-8", errors="ignore") as f:
        return f.read()[:200000]

def _ocr(path):
    try:
        import pytesseract
        from PIL import Image
        return pytesseract.image_to_string(Image.open(path))
    except Exception as e:
        return f"[image {os.path.basename(path)} — OCR unavailable ({e}); attach a text-based version for extraction]"

# --------------------------------------------------------------------------- #
#  Classifier
# --------------------------------------------------------------------------- #
CLASS_PATTERNS = [
    ("pacing",     r"pacing|scope\s*&?\s*sequence|unit\s+calendar|instructional days|week\s*\d"),
    ("ilp",        r"\bILP\b|instructional learning plan|learning plan"),
    ("perf_task",  r"performance task|rubric|exemplar|culminating task|PBA"),
    ("standards",  r"standard|GSE|CCSS|\b\d\.[A-Z]{1,3}\.\d"),
    ("roster",     r"roster|student name|diagnostic|score|tier\s*[123]|IEP|proficiency"),
    ("slides",     r"slide|deck|\.pptx"),
]

def classify(text: str, filename: str = "") -> str:
    hay = f"{filename}\n{text[:4000]}".lower()
    scores = {}
    for label, pat in CLASS_PATTERNS:
        scores[label] = len(re.findall(pat, hay, re.I))
    # filename hints
    fn = filename.lower()
    if fn.endswith(".pptx"): scores["slides"] = scores.get("slides", 0) + 3
    if fn.endswith((".xlsx", ".csv")): scores["roster"] = scores.get("roster", 0) + 2
    if "pacing" in fn: scores["pacing"] += 3
    if "ilp" in fn: scores["ilp"] += 3
    if "task" in fn: scores["perf_task"] += 2
    best = max(scores, key=scores.get) if scores else "other"
    return best if scores.get(best, 0) > 0 else "other"

def summarize_detection(text: str, label: str) -> str:
    stds = re.findall(r"\b\d+\.[A-Z]{1,3}\.?[A-Z]?\.?\d+[a-d]?\b", text)
    uniq = []
    for s in stds:
        if s not in uniq: uniq.append(s)
    bits = []
    if uniq: bits.append("standards " + ", ".join(uniq[:6]))
    days = re.findall(r"(\d+)\s*(?:instructional\s*)?days", text, re.I)
    if days: bits.append(f"{days[0]} days")
    return "; ".join(bits) or f"classified as {label}"
